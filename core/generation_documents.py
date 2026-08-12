"""
Génération de documents PDF à partir de texte/markdown.

Gratuit et local (WeasyPrint convertit du HTML/CSS en PDF, aucune clé API
requise) : contrairement à generation_images.py, cette fonctionnalité peut
rester active dès maintenant, elle ne coûte rien à faire tourner.

Flux : markdown -> HTML (lib `markdown`, déjà dans requirements.txt) ->
PDF (WeasyPrint) -> upload Supabase Storage -> URL publique renvoyée.

Prérequis Supabase à créer une fois, à la main, avant la première
utilisation (voir README_GENERATION.md) : un bucket public nommé
"generations", pas encore créé automatiquement par ce code.
"""

import html
import html
import logging
import uuid

import markdown as md_lib
from weasyprint import HTML

from api.auth import supabase
from core.conversion_pdf import conversion_disponible, convertir_en_pdf

BUCKET = "generations"

# Feuille de style minimale, volontairement sobre : Bourama pourra
# l'enrichir plus tard (logo Djiguignè, couleurs Maame) sans toucher à la
# logique de génération elle-même.
CSS_DE_BASE = """
@page { margin: 2.5cm; }
body { font-family: Helvetica, Arial, sans-serif; font-size: 11pt; line-height: 1.5; color: #1a1a1a; }
h1 { font-size: 20pt; margin-bottom: 0.3em; }
h2 { font-size: 15pt; margin-top: 1.2em; }
h3 { font-size: 12.5pt; margin-top: 1em; }
code { background: #f2f2f2; padding: 1px 4px; border-radius: 3px; }
pre { background: #f2f2f2; padding: 10px; border-radius: 5px; overflow-x: auto; }
table { border-collapse: collapse; width: 100%; }
td, th { border: 1px solid #ddd; padding: 6px 10px; }
"""


def generer_pdf_depuis_markdown(titre: str, contenu_markdown: str) -> str:
    """
    Convertit du markdown en PDF, l'upload dans Supabase Storage, renvoie
    l'URL publique.

    Lève une exception si WeasyPrint ou l'upload échoue -- à l'appelant
    (serveur MCP ou route REST) de transformer ça en message utilisateur
    clair, pas de logique de message d'erreur ici.
    """
    html_corps = md_lib.markdown(contenu_markdown, extensions=["tables", "fenced_code"])
    # CORRECTIF 2026-07-31 (audit sécurité/UX, même famille que le
    # correctif Excel : un titre parfaitement normal comme "Analyse < 100
    # unités & recommandations" suffisait à casser la mise en page du PDF
    # (titre tronqué/disparu), puisqu'il était inséré brut dans le HTML.
    # contenu_markdown n'est volontairement PAS touché ici : il passe déjà
    # par md_lib.markdown() qui produit du HTML propre à partir du
    # markdown écrit par le modèle.
    html_complet = f"""
    <html>
      <head><meta charset="utf-8"><style>{CSS_DE_BASE}</style></head>
      <body><h1>{html.escape(titre)}</h1>{html_corps}</body>
    </html>
    """

    pdf_bytes = HTML(string=html_complet).write_pdf()

    chemin = f"documents/{uuid.uuid4()}.pdf"
    try:
        supabase.storage.from_(BUCKET).upload(
            chemin, pdf_bytes, {"content-type": "application/pdf"}
        )
    except Exception as e:
        logging.error(f"ERREUR SUPABASE STORAGE (upload document {chemin}) : {e}")
        raise

    return supabase.storage.from_(BUCKET).get_public_url(chemin)


# --- Word / Excel / PowerPoint, ajoutés le 25/07 ---------------------------
#
# Contrairement au PDF ci-dessus (gratuit, local, WeasyPrint), ces formats
# ne sont pas "convertibles" directement en PDF sans dépendance lourde
# (LibreOffice) -- décision Bourama du 25/07 : passer par CloudConvert
# (service externe) plutôt que d'installer LibreOffice sur Railway. Voir
# core/conversion_pdf.py pour le détail et les compromis.
#
# Chaque fonction ci-dessous renvoie un dict {"url": ..., "url_apercu": ...}
# -- url_apercu est None si CLOUDCONVERT_API_KEY n'est pas configurée, ou si
# la conversion échoue pour une autre raison (quota dépassé, etc.) : dans
# ce cas le fichier original reste tout de même généré et téléchargeable,
# seul l'aperçu visuel manque. Ne JAMAIS faire échouer toute la génération
# à cause d'un souci d'aperçu, qui est une fonctionnalité secondaire par
# rapport au fichier lui-même.

MIME_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
MIME_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _nom_feuille_excel(titre: str) -> str:
    """
    CORRECTIF 2026-07-31 (audit sécurité/UX) : `titre` vient du modèle et
    était utilisé tel quel (juste tronqué à 31 caractères) comme nom de
    feuille -- Excel interdit pourtant les caractères : \\ / ? * [ ] dans
    un nom de feuille, ainsi qu'un nom vide ou commençant/finissant par
    une apostrophe. Un titre parfaitement normal ("Budget : prévisions
    2026", une date "01/08"...) suffisait à faire planter generer_xlsx.
    """
    interdits = ':\\/?*[]'
    nettoye = "".join(c for c in (titre or "") if c not in interdits).strip("'").strip()
    nettoye = nettoye[:31]
    return nettoye or "Feuille1"


def _uploader_avec_apercu(contenu_bytes: bytes, extension: str, content_type: str, nom_fichier: str) -> dict:
    """
    Upload le fichier original dans Supabase Storage, puis tente une
    conversion PDF pour l'aperçu (best-effort, voir note ci-dessus).
    Renvoie {"url": url_original, "url_apercu": url_pdf_ou_None}.
    """
    chemin = f"documents/{uuid.uuid4()}.{extension}"
    try:
        supabase.storage.from_(BUCKET).upload(chemin, contenu_bytes, {"content-type": content_type})
    except Exception as e:
        logging.error(f"ERREUR SUPABASE STORAGE (upload document {chemin}) : {e}")
        raise
    url = supabase.storage.from_(BUCKET).get_public_url(chemin)

    url_apercu = None
    if conversion_disponible():
        try:
            pdf_bytes = convertir_en_pdf(contenu_bytes, nom_fichier)
            chemin_apercu = f"documents/{uuid.uuid4()}_apercu.pdf"
            supabase.storage.from_(BUCKET).upload(chemin_apercu, pdf_bytes, {"content-type": "application/pdf"})
            url_apercu = supabase.storage.from_(BUCKET).get_public_url(chemin_apercu)
        except Exception as e:
            logging.warning(f"Aperçu PDF échoué pour {nom_fichier} (fichier original OK quand même) : {e}")

    return {"url": url, "url_apercu": url_apercu}


def generer_docx(titre: str, contenu_markdown: str) -> dict:
    """
    Génère un document Word simple à partir de markdown. Ne gère QUE
    titres (#, ##, ###) et paragraphes -- pas de gras/italique/listes
    imbriquées pour cette première version (contrairement au PDF qui
    passe par un vrai moteur HTML/CSS). À enrichir si Bourama en a besoin
    en usage réel.
    """
    import io
    import docx as docx_lib

    document = docx_lib.Document()
    document.add_heading(titre, level=0)

    for ligne in contenu_markdown.split("\n"):
        ligne_nettoyee = ligne.strip()
        if not ligne_nettoyee:
            continue
        if ligne_nettoyee.startswith("### "):
            document.add_heading(ligne_nettoyee[4:], level=3)
        elif ligne_nettoyee.startswith("## "):
            document.add_heading(ligne_nettoyee[3:], level=2)
        elif ligne_nettoyee.startswith("# "):
            document.add_heading(ligne_nettoyee[2:], level=1)
        else:
            document.add_paragraph(ligne_nettoyee)

    tampon = io.BytesIO()
    document.save(tampon)
    return _uploader_avec_apercu(tampon.getvalue(), "docx", MIME_DOCX, f"{titre}.docx")


def generer_xlsx(titre: str, en_tetes: list, lignes: list) -> dict:
    """
    Génère un classeur Excel à une feuille. `en_tetes` : liste de noms de
    colonnes. `lignes` : liste de listes (une sous-liste par ligne,
    valeurs dans le même ordre que en_tetes).
    """
    import io
    from openpyxl import Workbook

    classeur = Workbook()
    feuille = classeur.active
    feuille.title = _nom_feuille_excel(titre)  # voir _nom_feuille_excel : Excel interdit : \ / ? * [ ]

    feuille.append(en_tetes)
    for cellule in feuille[1]:
        cellule.font = cellule.font.copy(bold=True)
    for ligne in lignes:
        feuille.append(ligne)

    tampon = io.BytesIO()
    classeur.save(tampon)
    return _uploader_avec_apercu(tampon.getvalue(), "xlsx", MIME_XLSX, f"{titre}.xlsx")


def generer_pptx(titre: str, diapositives: list) -> dict:
    """
    Génère une présentation PowerPoint. `diapositives` : liste de dicts
    {"titre": str, "contenu": str} -- une diapositive titre+texte par
    élément, mise en page "Titre et contenu" standard (layout index 1 du
    template par défaut python-pptx).
    """
    import io
    from pptx import Presentation

    presentation = Presentation()

    diapo_titre = presentation.slides.add_slide(presentation.slide_layouts[0])
    diapo_titre.shapes.title.text = titre

    for diapo_info in diapositives:
        diapo = presentation.slides.add_slide(presentation.slide_layouts[1])
        diapo.shapes.title.text = diapo_info.get("titre", "")
        corps = diapo.placeholders[1].text_frame
        corps.text = diapo_info.get("contenu", "")

    tampon = io.BytesIO()
    presentation.save(tampon)
    return _uploader_avec_apercu(tampon.getvalue(), "pptx", MIME_PPTX, f"{titre}.pptx")
